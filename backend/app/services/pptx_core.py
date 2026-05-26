"""
PPT/PPTX 텍스트 추출기
- python-pptx: 기본 텍스트/표 추출
- EasyOCR + OpenCV: 이미지 내 텍스트 추출
- 읽기 순서 정렬 (top → left 기준)
"""

import io
import json
import re
from pathlib import Path
from typing import Any

import cv2
import easyocr
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

# ─────────────────────────────────────────────
# 상수
# ─────────────────────────────────────────────

# 읽기 순서 정렬 시 같은 행(row)으로 묶는 세로 허용 오차 (EMU 단위, 약 18pt)
ROW_TOLERANCE_EMU = 228600

# OCR 신뢰도 최소값 (0.0 ~ 1.0)
DEFAULT_CONFIDENCE_THRESHOLD = 0.3

# EasyOCR 지원 언어
OCR_LANGUAGES = ["ko", "en"]


# ─────────────────────────────────────────────
# OCR 리더 싱글톤 (초기화 비용 절감)
# ─────────────────────────────────────────────

_ocr_reader: easyocr.Reader | None = None


def get_ocr_reader() -> easyocr.Reader:
    global _ocr_reader
    if _ocr_reader is None:
        _ocr_reader = easyocr.Reader(OCR_LANGUAGES, gpu=False)
    return _ocr_reader


# ─────────────────────────────────────────────
# 이미지 전처리 (OpenCV)
# ─────────────────────────────────────────────

def preprocess_image_for_ocr(pil_image: Image.Image) -> np.ndarray:
    """
    OCR 인식률을 높이기 위한 이미지 전처리.
    흑백 변환 → 노이즈 제거 → 대비 향상 → 업스케일링 순으로 처리.
    """
    img = np.array(pil_image.convert("RGB"))
    img_bgr = cv2.cvtColor(img, cv2.COLOR_RGB2BGR)

    # 흑백 변환
    gray = cv2.cvtColor(img_bgr, cv2.COLOR_BGR2GRAY)

    # 해상도가 낮으면 2배 업스케일 (OCR 정확도 향상)
    h, w = gray.shape
    if w < 800:
        gray = cv2.resize(gray, (w * 2, h * 2), interpolation=cv2.INTER_CUBIC)

    # 가우시안 블러로 노이즈 제거
    gray = cv2.GaussianBlur(gray, (3, 3), 0)

    # CLAHE (Contrast Limited Adaptive Histogram Equalization) 로 대비 향상
    clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
    gray = clahe.apply(gray)

    # 적응형 이진화 (배경이 균일하지 않은 경우 효과적)
    binary = cv2.adaptiveThreshold(
        gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2
    )

    return binary


# ─────────────────────────────────────────────
# OCR 수행
# ─────────────────────────────────────────────

def run_ocr_on_image(
    pil_image: Image.Image,
    confidence_threshold: float = DEFAULT_CONFIDENCE_THRESHOLD,
    preprocess: bool = True,
) -> str:
    """
    PIL Image를 받아 EasyOCR로 텍스트를 추출한다.
    confidence_threshold 미만의 결과는 제거.
    """
    reader = get_ocr_reader()

    if preprocess:
        img_array = preprocess_image_for_ocr(pil_image)
    else:
        img_array = np.array(pil_image.convert("RGB"))

    results = reader.readtext(img_array, detail=1)

    lines: list[tuple[float, str]] = []
    for bbox, text, conf in results:
        if conf < confidence_threshold:
            continue
        # bbox의 y 중심값 기준으로 나중에 정렬
        y_center = (bbox[0][1] + bbox[2][1]) / 2
        lines.append((y_center, text.strip()))

    # 위 → 아래 순으로 정렬 후 합치기
    lines.sort(key=lambda t: t[0])
    return "\n".join(t for _, t in lines if t)


# ─────────────────────────────────────────────
# 텍스트 후처리
# ─────────────────────────────────────────────

def clean_text(text: str) -> str:
    """불필요한 공백·줄바꿈 정규화."""
    # 3개 이상 연속 줄바꿈 → 2개로
    text = re.sub(r"\n{3,}", "\n\n", text)
    # 줄 내 중복 공백 제거
    text = re.sub(r"[ \t]+", " ", text)
    return text.strip()


def deduplicate_texts(texts: list[str], existing: set[str]) -> list[str]:
    """
    이미 수집된 텍스트 집합(existing)과 겹치는 항목을 제거.
    OCR로 추출한 텍스트가 python-pptx 텍스트와 중복되는 경우를 막는다.
    """
    result = []
    for t in texts:
        normalized = re.sub(r"\s+", "", t)
        if normalized and normalized not in existing:
            existing.add(normalized)
            result.append(t)
    return result


# ─────────────────────────────────────────────
# 읽기 순서 정렬
# ─────────────────────────────────────────────

def sort_by_reading_order(
    blocks: list[dict], tolerance: int = ROW_TOLERANCE_EMU
) -> list[dict]:
    """
    블록 목록을 사람의 읽기 순서(위→아래, 왼→오른)로 정렬.
    top 값이 tolerance 이내면 같은 행으로 간주하고 left로 2차 정렬.
    """
    if not blocks:
        return blocks

    sorted_blocks = sorted(blocks, key=lambda b: (b.get("top", 0), b.get("left", 0)))

    # 행 그룹화
    rows: list[list[dict]] = []
    current_row: list[dict] = [sorted_blocks[0]]
    row_top = sorted_blocks[0].get("top", 0)

    for block in sorted_blocks[1:]:
        top = block.get("top", 0)
        if abs(top - row_top) <= tolerance:
            current_row.append(block)
        else:
            rows.append(sorted(current_row, key=lambda b: b.get("left", 0)))
            current_row = [block]
            row_top = top

    rows.append(sorted(current_row, key=lambda b: b.get("left", 0)))

    return [block for row in rows for block in row]


# ─────────────────────────────────────────────
# 표(Table) 추출
# ─────────────────────────────────────────────

def extract_table(shape) -> list[list[str]]:
    """
    테이블 shape에서 행/열 구조를 유지하며 텍스트 추출.
    병합 셀(span_height > 1 또는 span_width > 1)은 첫 번째 셀에만 텍스트를 넣고
    나머지는 빈 문자열 "<merged>"로 표시.
    """
    table = shape.table
    rows_data: list[list[str]] = []

    # 병합 셀 추적용 집합
    merged_positions: set[tuple[int, int]] = set()

    for r_idx, row in enumerate(table.rows):
        row_data: list[str] = []
        for c_idx, cell in enumerate(row.cells):
            if (r_idx, c_idx) in merged_positions:
                row_data.append("<merged>")
                continue

            text = cell.text.strip()
            row_data.append(text)

            # 병합된 셀 범위 기록
            sp = cell._tc
            row_span = int(sp.get("rowSpan", 1) or 1)
            col_span = int(sp.get("gridSpan", 1) or 1)
            for dr in range(row_span):
                for dc in range(col_span):
                    if dr == 0 and dc == 0:
                        continue
                    merged_positions.add((r_idx + dr, c_idx + dc))

        rows_data.append(row_data)

    return rows_data


# ─────────────────────────────────────────────
# Shape 텍스트 추출 (재귀적으로 그룹/SmartArt 처리)
# ─────────────────────────────────────────────

def extract_text_blocks_from_shape(
    shape,
    collected_texts: set[str],
    confidence_threshold: float,
) -> list[dict]:
    """
    단일 shape에서 텍스트 블록 목록을 반환.
    그룹(Group) 또는 SmartArt shape는 재귀 탐색.
    이미지 shape는 OCR 수행.
    """
    blocks: list[dict] = []

    # 위치 정보 (EMU 단위)
    try:
        top = shape.top or 0
        left = shape.left or 0
    except Exception:
        top, left = 0, 0

    shape_type = shape.shape_type

    # ── 그룹 shape 재귀 탐색 ──────────────────────
    if shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            blocks.extend(
                extract_text_blocks_from_shape(child, collected_texts, confidence_threshold)
            )
        return blocks

    # ── 표 ───────────────────────────────────────
    if shape.has_table:
        # 표는 tables 섹션에서 별도 처리하므로 여기선 스킵
        return blocks

    # ── 일반 텍스트 프레임 ────────────────────────
    if shape.has_text_frame:
        text_parts = []
        for para in shape.text_frame.paragraphs:
            para_text = "".join(run.text for run in para.runs).strip()
            if para_text:
                text_parts.append(para_text)
        full_text = clean_text("\n".join(text_parts))

        if full_text:
            normalized = re.sub(r"\s+", "", full_text)
            if normalized not in collected_texts:
                collected_texts.add(normalized)
                blocks.append(
                    {
                        "type": "text",
                        "content": full_text,
                        "top": top,
                        "left": left,
                    }
                )
        return blocks

    # ── 이미지 → OCR ──────────────────────────────
    if shape_type in (
        MSO_SHAPE_TYPE.PICTURE,
        MSO_SHAPE_TYPE.LINKED_PICTURE,
    ):
        try:
            image_blob = shape.image.blob
            pil_image = Image.open(io.BytesIO(image_blob))
            ocr_text = run_ocr_on_image(pil_image, confidence_threshold)
            ocr_text = clean_text(ocr_text)

            if ocr_text:
                deduped = deduplicate_texts([ocr_text], collected_texts)
                for t in deduped:
                    blocks.append(
                        {
                            "type": "ocr",
                            "content": t,
                            "top": top,
                            "left": left,
                        }
                    )
        except Exception as e:
            blocks.append(
                {
                    "type": "ocr_error",
                    "content": f"[OCR 실패: {e}]",
                    "top": top,
                    "left": left,
                }
            )
        return blocks

    # ── 도형 내부 텍스트 (그림 아닌 shape) ─────────
    # MSO_SHAPE_TYPE.AUTO_SHAPE, LINE, FREEFORM 등
    try:
        if hasattr(shape, "text") and shape.text.strip():
            full_text = clean_text(shape.text)
            normalized = re.sub(r"\s+", "", full_text)
            if normalized not in collected_texts:
                collected_texts.add(normalized)
                blocks.append(
                    {
                        "type": "shape_text",
                        "content": full_text,
                        "top": top,
                        "left": left,
                    }
                )
    except Exception:
        pass

    return blocks


# ─────────────────────────────────────────────
# 슬라이드 처리
# ─────────────────────────────────────────────

def process_slide(
    slide,
    slide_number: int,
    confidence_threshold: float = DEFAULT_CONFIDENCE_THRESHOLD,
) -> dict[str, Any]:
    """
    슬라이드 하나를 처리하여 JSON 직렬화 가능한 딕셔너리 반환.
    text_blocks: 정렬된 텍스트/OCR 블록 목록
    tables: 2D 문자열 배열 목록
    """
    collected_texts: set[str] = set()
    all_text_blocks: list[dict] = []
    tables: list[list[list[str]]] = []

    for shape in slide.shapes:
        # 표는 별도 수집
        if shape.has_table:
            tables.append(extract_table(shape))
            # 표 셀 텍스트를 collected_texts에 등록해 OCR 중복 방지
            for row in tables[-1]:
                for cell_text in row:
                    if cell_text and cell_text != "<merged>":
                        collected_texts.add(re.sub(r"\s+", "", cell_text))
            continue

        blocks = extract_text_blocks_from_shape(shape, collected_texts, confidence_threshold)
        all_text_blocks.extend(blocks)

    # 읽기 순서 정렬 후 위치 정보 제거 (출력 불필요)
    sorted_blocks = sort_by_reading_order(all_text_blocks)
    clean_blocks = [
        {"type": b["type"], "content": b["content"]} for b in sorted_blocks
    ]

    return {
        "slide_number": slide_number,
        "text_blocks": clean_blocks,
        "tables": tables,
    }


# ─────────────────────────────────────────────
# 메인 추출 함수
# ─────────────────────────────────────────────

def extract_text_from_pptx(
    file_path: str | Path,
    confidence_threshold: float = DEFAULT_CONFIDENCE_THRESHOLD,
    slide_range: tuple[int, int] | None = None,
) -> list[dict[str, Any]]:
    """
    PPTX 파일에서 전체 슬라이드 텍스트를 추출하여 반환.

    Args:
        file_path: .pptx 파일 경로
        confidence_threshold: OCR 신뢰도 최소값 (기본 0.3)
        slide_range: 처리할 슬라이드 범위 (1-based). None이면 전체.
                     예: (2, 5) → 2~5번 슬라이드만 처리

    Returns:
        슬라이드별 딕셔너리 목록
    """
    prs = Presentation(str(file_path))
    results: list[dict] = []

    total = len(prs.slides)
    start, end = 1, total
    if slide_range:
        start = max(1, slide_range[0])
        end = min(total, slide_range[1])

    for idx, slide in enumerate(prs.slides, start=1):
        if idx < start or idx > end:
            continue
        print(f"  슬라이드 {idx}/{total} 처리 중...")
        slide_data = process_slide(slide, idx, confidence_threshold)
        results.append(slide_data)

    return results


# ─────────────────────────────────────────────
# TXT 저장 (result 폴더, 슬라이드별 개별 파일)
# ─────────────────────────────────────────────

def format_slide_as_text(slide_data: dict[str, Any]) -> str:
    """
    슬라이드 딕셔너리를 사람이 읽기 좋은 텍스트로 변환.
    text_blocks → 본문, tables → 격자 형태로 출력.
    """
    lines: list[str] = []
    slide_num = slide_data["slide_number"]

    lines.append(f"{'=' * 50}")
    lines.append(f"[ 슬라이드 {slide_num} ]")
    lines.append(f"{'=' * 50}")

    # ── 텍스트 블록 ──────────────────────────────
    for block in slide_data.get("text_blocks", []):
        tag = {"text": "", "shape_text": "[도형] ", "ocr": "[OCR] ", "ocr_error": ""}.get(
            block["type"], ""
        )
        lines.append(f"{tag}{block['content']}")

    # ── 표 ───────────────────────────────────────
    for t_idx, table in enumerate(slide_data.get("tables", []), start=1):
        lines.append(f"\n[표 {t_idx}]")
        if not table:
            continue
        # 열 너비를 내용 최대 길이로 맞춤
        col_widths = [
            max(len(str(row[c])) for row in table if c < len(row))
            for c in range(max(len(row) for row in table))
        ]
        for row in table:
            cells = [str(row[c]) if c < len(row) else "" for c in range(len(col_widths))]
            padded = [cell.ljust(col_widths[i]) for i, cell in enumerate(cells)]
            lines.append("| " + " | ".join(padded) + " |")

    lines.append("")  # 슬라이드 간 빈 줄
    return "\n".join(lines)
